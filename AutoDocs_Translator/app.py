import streamlit as st
import os
import io
import pypdf
from docx import Document
import openpyxl
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
import ollama
from langchain_community.embeddings import OllamaEmbeddings
from langchain_chroma import Chroma
from langchain_core.documents import Document as LcDocument

st.set_page_config(page_title="AutoDocs Translator & AI Assistant", layout="wide", page_icon="logo.png")

# CSS to make the UI look more premium
st.markdown("""
<style>
    .reportview-container {
        background: #0E1117;
    }
    .stButton>button {
        border-radius: 8px;
        transition: all 0.3s ease;
    }
    .stButton>button:hover {
        transform: scale(1.02);
    }
    .translation-box {
        background-color: #1E2329;
        padding: 15px;
        border-radius: 8px;
        border: 1px solid #333;
        height: 100%;
        overflow-y: auto;
    }
</style>
""", unsafe_allow_html=True)

# Sidebar Navigation
st.sidebar.image("logo.png", use_container_width=True)
st.sidebar.title("🛠️ Chế độ hoạt động")
mode = st.sidebar.radio("Chọn chức năng:", ["📝 Dịch Thuật Tài Liệu (Word, PDF...)", "🤖 AI Chatbot (Học & Hỏi đáp)"])
st.sidebar.markdown("---")
st.sidebar.info("Công cụ nội bộ Air-Gapped. Chạy hoàn toàn Offline trên máy tính cá nhân.")

# ==================== TEXT EXTRACTION ====================
def extract_text_from_pdf(file_bytes):
    text = ""
    reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    for page in reader.pages:
        if page.extract_text():
            text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_xlsx(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text += f"\n--- Sheet: {sheet} ---\n"
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                text += row_text + "\n"
    return text

def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def process_file(uploaded_file):
    name = uploaded_file.name.lower()
    file_bytes = uploaded_file.read()
    
    if name.endswith('.pdf'):
        return extract_text_from_pdf(file_bytes)
    elif name.endswith('.docx'):
        return extract_text_from_docx(file_bytes)
    elif name.endswith('.xlsx'):
        return extract_text_from_xlsx(file_bytes)
    elif name.endswith('.pptx'):
        return extract_text_from_pptx(file_bytes)
    elif name.endswith('.txt'):
        return file_bytes.decode('utf-8')
    else:
        st.error("Định dạng file không được hỗ trợ!")
        return ""

def create_docx(translated_texts):
    doc = Document()
    doc.add_heading("Bản Dịch Tài Liệu", 0)
    for text in translated_texts:
        doc.add_paragraph(text)
    
    # Save to BytesIO
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# ==================== OLLAMA LLM ====================
def translate_chunk(chunk_text, context_prompt=""):
    system_prompt = (
        "Bạn là một kỹ sư hệ thống điện/điện tử ô tô (Automotive E/E Engineer). "
        "Nhiệm vụ của bạn là dịch tài liệu kỹ thuật sau từ Tiếng Anh sang Tiếng Việt. "
        "Yêu cầu: "
        "1. Giữ nguyên các từ viết tắt chuyên ngành như CAN, LIN, ECU, BCM, VCU, ADAS, Gateway, v.v. "
        "2. Văn phong chuẩn xác, rõ ràng, dễ hiểu cho kỹ sư người Việt. "
        "3. Tuyệt đối không tự bịa thêm thông tin không có trong văn bản gốc. "
        f"{context_prompt}"
    )
    
    try:
        response = ollama.chat(model='qwen2', messages=[
            {'role': 'system', 'content': system_prompt},
            {'role': 'user', 'content': f"Hãy dịch đoạn văn bản sau:\n\n{chunk_text}"}
        ])
        return response['message']['content']
    except Exception as e:
        return f"[Lỗi Ollama: {str(e)}]"


# ==================== MODE 1: TRANSLATOR ====================
if "Dịch Thuật" in mode:
    st.title("🛡️ AutoDocs Translator")
    st.write("Dịch thuật tài liệu kỹ thuật Mật. Đầu ra hỗ trợ xuất file Word (.docx) chuyên nghiệp.")
    
    uploaded_file = st.file_uploader("Tải tài liệu mật lên (Hỗ trợ: PDF, DOCX, XLSX, PPTX, TXT)", type=['pdf', 'docx', 'xlsx', 'pptx', 'txt'])

    if uploaded_file is not None:
        with st.spinner("Đang trích xuất văn bản..."):
            raw_text = process_file(uploaded_file)
            
        if raw_text.strip():
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1500,
                chunk_overlap=200,
                separators=["\n\n", "\n", ".", " ", ""]
            )
            chunks = text_splitter.split_text(raw_text)
            
            st.success(f"Tài liệu đã được trích xuất thành công và chia thành {len(chunks)} đoạn để dịch.")
            
            if st.button("🚀 Bắt đầu dịch toàn bộ", type="primary"):
                st.session_state['translating'] = True
                st.session_state['translated_chunks'] = []
            
            if st.session_state.get('translating'):
                translated_results = []
                
                # Show side-by-side
                st.markdown("### Quá trình dịch thuật:")
                
                for i, chunk in enumerate(chunks):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.markdown(f"**Bản Gốc (Đoạn {i+1})**")
                        st.markdown(f"<div class='translation-box'>{chunk}</div>", unsafe_allow_html=True)
                        
                    with col2:
                        st.markdown(f"**Bản Dịch (Đoạn {i+1})**")
                        with st.spinner("AI đang dịch..."):
                            translated = translate_chunk(chunk)
                            st.markdown(f"<div class='translation-box'>{translated}</div>", unsafe_allow_html=True)
                            translated_results.append(translated)
                    
                    st.write("---")
                
                st.success("Hoàn thành dịch thuật! Bạn có thể tải file Word về bên dưới.")
                
                # Download as DOCX
                docx_buffer = create_docx(translated_results)
                st.download_button(
                    label="📥 Tải xuống Bản Dịch (.docx)",
                    data=docx_buffer,
                    file_name=f"{uploaded_file.name}_translated.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )

# ==================== MODE 2: RAG CHATBOT ====================
elif "AI Chatbot" in mode:
    st.title("🤖 AI Teacher (Hỏi Đáp Tài Liệu)")
    st.write("Cung cấp tài liệu Mật, AI sẽ 'học' và trả lời mọi câu hỏi của bạn dựa trên tài liệu đó.")
    
    rag_file = st.file_uploader("Tải tài liệu muốn AI học (Hỗ trợ: PDF, DOCX, TXT...)", type=['pdf', 'docx', 'xlsx', 'pptx', 'txt'], key="rag_file")
    
    if rag_file is not None:
        if st.button("🧠 Cho AI học tài liệu này"):
            with st.spinner("Đang trích xuất và nạp dữ liệu vào Não AI (Vector Database)..."):
                raw_text = process_file(rag_file)
                text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
                chunks = text_splitter.split_text(raw_text)
                
                # Convert to Langchain Documents
                documents = [LcDocument(page_content=t) for t in chunks]
                
                # Setup Embeddings and Vector DB
                embeddings = OllamaEmbeddings(model="nomic-embed-text")
                
                # Clear old DB if exists to avoid memory leak
                if "vector_store" in st.session_state:
                    del st.session_state["vector_store"]
                
                vector_store = Chroma.from_documents(documents, embeddings)
                st.session_state["vector_store"] = vector_store
                st.session_state["chat_history"] = []
                
                st.success("🎉 AI đã học xong! Bạn có thể hỏi đáp ngay bây giờ.")
                
    if "vector_store" in st.session_state:
        st.markdown("---")
        st.subheader("Trò chuyện với Tài Liệu")
        
        # Display chat history
        for msg in st.session_state.get("chat_history", []):
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])
                
        # Chat Input
        if prompt := st.chat_input("Hỏi gì đó về tài liệu này (Ví dụ: Quy trình nổ máy như nào?)"):
            # Add user message
            st.session_state["chat_history"].append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)
                
            # RAG Retrieval
            with st.spinner("Đang tìm kiếm thông tin trong tài liệu..."):
                retriever = st.session_state["vector_store"].as_retriever(search_kwargs={"k": 3})
                relevant_docs = retriever.invoke(prompt)
                context = "\n\n".join([doc.page_content for doc in relevant_docs])
                
                system_prompt = (
                    "Bạn là một người Thầy dạy kèm, chuyên gia kỹ thuật Ô tô. "
                    "Hãy sử dụng NỘI DUNG TÀI LIỆU được cung cấp dưới đây để trả lời câu hỏi của người dùng bằng Tiếng Việt. "
                    "Nếu thông tin KHÔNG có trong tài liệu, hãy nói rõ là tài liệu không đề cập đến. "
                    f"NỘI DUNG TÀI LIỆU:\n{context}"
                )
                
                try:
                    response = ollama.chat(model='qwen2', messages=[
                        {'role': 'system', 'content': system_prompt},
                        {'role': 'user', 'content': prompt}
                    ])
                    ai_reply = response['message']['content']
                except Exception as e:
                    ai_reply = f"Lỗi AI: {e}"
            
            # Display AI message
            with st.chat_message("assistant"):
                st.markdown(ai_reply)
                with st.expander("Nguồn trích dẫn (Click để xem)"):
                    st.write(context)
            
            st.session_state["chat_history"].append({"role": "assistant", "content": ai_reply})
